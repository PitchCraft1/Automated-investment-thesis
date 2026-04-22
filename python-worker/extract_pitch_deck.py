import io
import json
import os
import shutil
import subprocess
import sys
import tempfile

from pptx import Presentation

try:
    import pytesseract
    from PIL import Image
except ImportError:
    pytesseract = None
    Image = None


SECTION_KEYWORDS = {
    "problem_statement": ["problem", "pain", "challenge", "friction", "issue"],
    "solution_product": ["solution", "product", "platform", "workflow", "automation"],
    "market_size": ["tam", "sam", "som", "market", "opportunity"],
    "business_model": ["pricing", "revenue", "business model", "subscription"],
    "competitive_landscape": ["competitor", "competition", "advantage", "differentiation"],
    "team": ["founder", "team", "ceo", "cto", "advisor"],
    "financial_projections": ["forecast", "projection", "financial", "margin", "runway"],
    "traction_milestones": ["traction", "milestone", "growth", "users", "pilot", "revenue"],
    "funding_ask": ["funding", "ask", "raise", "investment", "runway"],
}


def detect_slide_types(text):
    lowered = text.lower()
    detected = []
    for key, keywords in SECTION_KEYWORDS.items():
        if any(keyword in lowered for keyword in keywords):
            detected.append(key)
    return detected


def convert_ppt_to_pptx(file_path):
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        raise RuntimeError("LibreOffice is required for .ppt conversion but was not found in PATH.")

    output_dir = tempfile.mkdtemp(prefix="ppt-convert-")
    subprocess.run(
        [soffice, "--headless", "--convert-to", "pptx", "--outdir", output_dir, file_path],
        check=True,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
    )
    converted_name = os.path.splitext(os.path.basename(file_path))[0] + ".pptx"
    converted_path = os.path.join(output_dir, converted_name)
    if not os.path.exists(converted_path):
        raise RuntimeError("PPT conversion completed but no PPTX output was created.")
    return converted_path


def extract_text_from_shape(shape):
    text_parts = []
    if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            paragraph_text = " ".join(run.text.strip() for run in paragraph.runs if run.text and run.text.strip())
            if paragraph_text:
                text_parts.append(paragraph_text)
    return text_parts


def extract_ocr_from_shape(shape):
    if pytesseract is None or Image is None:
        return []

    image = getattr(shape, "image", None)
    if image is None:
        return []

    try:
        picture = Image.open(io.BytesIO(image.blob))
        text = pytesseract.image_to_string(picture).strip()
        return [text] if text else []
    except Exception:
        return []


def load_presentation(file_path):
    if file_path.lower().endswith(".pptx"):
        return file_path
    if file_path.lower().endswith(".ppt"):
        return convert_ppt_to_pptx(file_path)
    raise RuntimeError("Unsupported file type. Only .ppt and .pptx are accepted.")


def extract_deck(file_path):
    resolved_path = load_presentation(file_path)
    presentation = Presentation(resolved_path)
    slides = []
    warnings = []

    if pytesseract is None:
        warnings.append("pytesseract and Pillow are not installed, so OCR was skipped for image-based slides.")

    for slide in presentation.slides:
        tokens = []
        for shape in slide.shapes:
            tokens.extend(extract_text_from_shape(shape))
            tokens.extend(extract_ocr_from_shape(shape))

        cleaned_tokens = [token.strip() for token in tokens if token and token.strip()]
        title = cleaned_tokens[0] if cleaned_tokens else ""
        slides.append({
            "title": title,
            "text": " ".join(cleaned_tokens),
            "tokens": cleaned_tokens,
        })

    full_text = "\n".join(slide["text"] for slide in slides)
    detected_slide_types = sorted(set(sum([detect_slide_types(slide["text"]) for slide in slides], [])))

    return {
        "slideCount": len(slides),
        "slides": slides,
        "fullText": full_text,
        "detectedSlideTypes": detected_slide_types,
        "warnings": warnings,
    }


def main():
    if len(sys.argv) < 2:
        raise SystemExit("File path is required.")

    file_path = sys.argv[1]
    if not os.path.exists(file_path):
        raise SystemExit("Pitch deck file does not exist.")

    print(json.dumps(extract_deck(file_path)))


if __name__ == "__main__":
    main()
